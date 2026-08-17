import asyncio
import io
import logging
import base64
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from anthropic import AsyncAnthropic
from app.core.config import get_settings
from app.core.supabase import get_supabase
from app.core.async_db import run_db_operation, db_storage_download
from app.services.analysis_service import AnalysisService

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Maximum time allowed for document processing (extraction + analysis) (10 minutes)
DOCUMENT_PROCESSING_TIMEOUT = 600
# Pages with fewer characters than this are treated as scanned (image-only) pages
SCANNED_PAGE_CHAR_THRESHOLD = 50
# Claude Vision model for OCR on images and scanned PDF pages
VISION_MODEL = "claude-sonnet-4-6"


class ExtractionService:
    def __init__(self):
        self.supabase = get_supabase()
        self.bucket_name = "course-materials"
        self.analysis_service = AnalysisService()
        settings = get_settings()
        self.vision_client = AsyncAnthropic(api_key=settings.ANTHROPIC_API_KEY)

    async def process_document(self, document_id: str):
        """
        Process a document through the extraction pipeline:
        1. Download from storage
        2. Extract text
        3. Analyze with Claude (extract topics/concepts)

        Document is marked 'completed' only once concepts have actually been
        extracted. Quiz question generation is a separate, user-triggered flow
        (on_demand_quiz_service) — this no longer auto-generates a quiz on
        upload; that quiz bank was never shown to students and existed only
        to gate a "materials ready" signal, which is now derived for free
        from this document reaching 'completed'.

        All database operations are async to prevent blocking the event loop.
        """
        try:
            await asyncio.wait_for(
                self._process_document_internal(document_id),
                timeout=DOCUMENT_PROCESSING_TIMEOUT
            )
        except asyncio.TimeoutError:
            logger.error(f"Document {document_id}: Processing timed out after {DOCUMENT_PROCESSING_TIMEOUT}s")
            await self._update_status_with_error(
                document_id,
                "failed",
                "Processing timed out. The document may be too large or complex. Please try a smaller document."
            )
        except Exception as e:
            logger.exception(f"Document {document_id}: Unexpected error during processing")
            await self._update_status_with_error(
                document_id,
                "failed",
                "An unexpected error occurred during processing. Please try again."
            )

    async def _process_document_internal(self, document_id: str):
        """Internal processing logic: extraction + analysis."""
        try:
            logger.info(f"Starting processing for document {document_id}")

            # 1. Get document metadata (ASYNC)
            doc_response = await run_db_operation(
                lambda: self.supabase.table("documents").select("*").eq("id", document_id).single().execute()
            )

            if not doc_response.data:
                logger.error(f"Document {document_id} not found in database")
                await self._update_status_with_error(document_id, "failed", "Document not found in database")
                return

            doc = doc_response.data

            # Type check the document data
            if not isinstance(doc, dict):
                logger.error(f"Invalid document data format for {document_id}")
                await self._update_status_with_error(document_id, "failed", "Invalid document data format")
                return

            # 2. Update status to processing (ASYNC)
            await self._update_status(document_id, "processing")
            logger.info(f"Document {document_id}: Status updated to 'processing'")

            # 3. Download file (ASYNC)
            file_path = doc.get("file_path")
            if not file_path or not isinstance(file_path, str):
                raise ValueError("Document is missing file path")

            logger.info(f"Document {document_id}: Downloading file from storage...")
            try:
                file_content = await db_storage_download(self.supabase, self.bucket_name, file_path)
            except Exception as e:
                raise ValueError(f"Failed to download file from storage: {str(e)}")

            # 4. Extract text (CPU-bound, run in executor to not block)
            file_type = doc.get("file_type")
            if not file_type or not isinstance(file_type, str):
                raise ValueError("Document is missing file type")

            logger.info(f"Document {document_id}: Extracting text from {file_type}...")
            try:
                extracted_text = await self._extract_text_async(file_content, file_type)
            except ValueError as e:
                raise ValueError(f"Text extraction failed: {str(e)}")

            if not extracted_text or len(extracted_text.strip()) < 50:
                raise ValueError(
                    "This document has too little text to work with. "
                    "Try uploading a more complete version."
                )

            # Save extracted text immediately so it's not lost if analysis fails (ASYNC)
            await run_db_operation(
                lambda: self.supabase.table("documents").update({
                    "extracted_text": extracted_text
                }).eq("id", document_id).execute()
            )
            logger.info(f"Document {document_id}: Extracted {len(extracted_text)} characters of text")

            # 5. Extract Structure (Topics & Concepts)
            logger.info(f"Document {document_id}: Starting AI analysis...")
            try:
                await self.analysis_service.analyze_document(document_id, extracted_text)
            except Exception as e:
                logger.error(f"Document {document_id}: Analysis failed - {str(e)}")
                raise ValueError(
                    "We had trouble reading this document's content. "
                    "Please try again — or try uploading a different file format."
                )

            # Verify concepts were created (ASYNC)
            concepts_count = await self._count_document_concepts(document_id)
            if concepts_count == 0:
                raise ValueError(
                    "This document doesn't have enough learning material for PassAI to work with. "
                    "Try uploading a more detailed version."
                )

            logger.info(f"Document {document_id}: Analysis complete - {concepts_count} concepts extracted")

            # 6. Mark document as completed (ASYNC)
            await run_db_operation(
                lambda: self.supabase.table("documents").update({
                    "status": "completed",
                    "error_message": None
                }).eq("id", document_id).execute()
            )
            logger.info(f"Document {document_id}: Status updated to 'completed'")

        except ValueError as e:
            # User-friendly errors (validation, unsupported file types, etc.)
            error_message = str(e)
            logger.error(f"Document {document_id}: Processing failed - {error_message}")
            await self._update_status_with_error(document_id, "failed", error_message)

    async def _count_document_concepts(self, document_id: str) -> int:
        """Count the number of concepts extracted for a document (ASYNC)."""
        try:
            # Get topics for document
            topics_response = await run_db_operation(
                lambda: self.supabase.table("topics").select("id").eq("document_id", document_id).execute()
            )
            if not topics_response.data:
                return 0

            topic_ids = [t["id"] for t in topics_response.data]

            # Count concepts for those topics
            concepts_response = await run_db_operation(
                lambda: self.supabase.table("concepts").select("id", count="exact").in_("topic_id", topic_ids).execute()
            )
            return concepts_response.count if concepts_response.count else 0
        except Exception:
            return 0

    async def _update_status_with_error(self, document_id: str, status: str, error_message: str):
        """Update document status and error message (ASYNC)."""
        try:
            await run_db_operation(
                lambda: self.supabase.table("documents").update({
                    "status": status,
                    "error_message": error_message
                }).eq("id", document_id).execute()
            )
        except Exception as e:
            logger.error(f"Failed to update document {document_id} status: {e}")

    async def _update_status(self, document_id: str, status: str):
        """Update document status (ASYNC)."""
        await run_db_operation(
            lambda: self.supabase.table("documents").update({"status": status}).eq("id", document_id).execute()
        )

    async def _extract_text_async(self, file_content: bytes, file_type: str) -> str:
        """
        Async text extraction dispatcher. Routes to the appropriate extractor based on file type:
        - Images (JPEG, PNG, GIF, WEBP): Claude Vision OCR
        - PDFs (native text + scanned): _parse_pdf (hybrid)
        - DOCX, PPTX, plain text: _extract_text_sync in executor (CPU-bound)
        """
        ft = file_type.lower()

        if any(t in ft for t in ("jpeg", "jpg", "png", "gif", "webp", "image")):
            media_type = self._normalize_media_type(file_type)
            logger.info(f"Extracting text from image via Vision OCR ({media_type})")
            text = await self._ocr_image_bytes(file_content, media_type)
            if not text:
                raise ValueError(
                    "Could not extract text from the image. "
                    "Please ensure the image contains readable text."
                )
            return text

        if "pdf" in ft:
            logger.info("Extracting text from PDF (with scanned-page detection)")
            return await self._parse_pdf(file_content)

        # DOCX, PPTX, plain text — synchronous, run in executor
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            None,
            lambda: self._extract_text_sync(file_content, file_type)
        )

    async def _parse_pdf(self, file_content: bytes) -> str:
        """
        Extract text from a PDF using PyMuPDF.
        Pages with fewer than SCANNED_PAGE_CHAR_THRESHOLD characters of native text
        are treated as scanned and processed with Claude Vision OCR.
        """
        loop = asyncio.get_event_loop()
        pdf_doc = await loop.run_in_executor(
            None, lambda: fitz.open(stream=file_content, filetype="pdf")
        )
        try:
            pages_text = []
            for page_num in range(len(pdf_doc)):
                page = pdf_doc[page_num]
                native_text = await loop.run_in_executor(None, page.get_text)

                if len(native_text.strip()) >= SCANNED_PAGE_CHAR_THRESHOLD:
                    pages_text.append(native_text.strip())
                else:
                    logger.info(
                        f"Page {page_num + 1}: native text too short "
                        f"({len(native_text.strip())} chars), using Vision OCR"
                    )
                    matrix = fitz.Matrix(2, 2)
                    pixmap = await loop.run_in_executor(
                        None, lambda p=page, m=matrix: p.get_pixmap(matrix=m)
                    )
                    png_bytes = await loop.run_in_executor(None, pixmap.tobytes, "png")
                    ocr_text = await self._ocr_image_bytes(png_bytes, "image/png")
                    if ocr_text:
                        pages_text.append(ocr_text)
                    else:
                        logger.warning(f"Page {page_num + 1}: Vision OCR returned empty text")

            return "\n\n".join(pages_text).strip()
        finally:
            pdf_doc.close()

    async def _ocr_image_bytes(self, image_bytes: bytes, media_type: str) -> str:
        """
        Send image bytes to Claude Vision and return extracted text.
        Returns empty string on failure so callers can decide how to handle it.
        """
        try:
            encoded = base64.standard_b64encode(image_bytes).decode("utf-8")
            response = await self.vision_client.messages.create(
                model=VISION_MODEL,
                max_tokens=4096,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "image",
                                "source": {
                                    "type": "base64",
                                    "media_type": media_type,
                                    "data": encoded,
                                },
                            },
                            {
                                "type": "text",
                                "text": (
                                    "Extract all text from this image exactly as it appears. "
                                    "Preserve paragraph breaks. "
                                    "Output only the extracted text, no commentary."
                                ),
                            },
                        ],
                    }
                ],
            )
            return response.content[0].text.strip()
        except Exception as e:
            logger.warning(f"Vision OCR failed: {e}")
            return ""

    @staticmethod
    def _normalize_media_type(file_type: str) -> str:
        """Map a file_type string to a Claude Vision-compatible media type."""
        ft = file_type.lower()
        if "jpeg" in ft or "jpg" in ft:
            return "image/jpeg"
        if "png" in ft:
            return "image/png"
        if "gif" in ft:
            return "image/gif"
        if "webp" in ft:
            return "image/webp"
        raise ValueError(f"Unsupported image media type: {file_type}")

    def _extract_text_sync(self, file_content: bytes, file_type: str) -> str:
        """
        Synchronous text extraction for DOCX, PPTX, and plain text.
        Intended to run in an executor (CPU-bound).
        """
        text = ""
        file_stream = io.BytesIO(file_content)

        if "word" in file_type or "docx" in file_type:
            # application/vnd.openxmlformats-officedocument.wordprocessingml.document
            doc = docx.Document(file_stream)

            # Extract from paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        text += " | ".join(row_text) + "\n"

        elif "powerpoint" in file_type or "pptx" in file_type or "presentation" in file_type:
            # application/vnd.openxmlformats-officedocument.presentationml.presentation
            prs = Presentation(file_stream)

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []

                for shape in slide.shapes:
                    # Extract text from text frames
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = ""
                            for run in paragraph.runs:
                                if run.text:
                                    para_text += run.text
                            if para_text.strip():
                                slide_text.append(para_text.strip())

                    # Extract text from tables in slides
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))

                if slide_text:
                    text += f"--- Slide {slide_num} ---\n"
                    text += "\n".join(slide_text) + "\n\n"

        elif "text" in file_type:  # text/plain
            text = file_content.decode('utf-8')

        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        return text.strip()
